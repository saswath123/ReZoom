from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
from datetime import datetime
import re

def remove_unicode(text):
    if not text:
        return ""

    return text.encode("ascii", "ignore").decode("ascii")

def safe_str(value):
    """Convert None to empty string safely"""
    return str(value) if value is not None else ""

def create_professional_avatar(slide, left, top, size, name, gender="neutral"):
    """Create a professional avatar with initials"""
    
    # Create circular shape for avatar background
    avatar = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left),
        Inches(top),
        Inches(size),
        Inches(size)
    )
    
    # Professional color based on gender (subtle, not cartoonish)
    if gender == "male":
        bg_color = RGBColor(52, 73, 94)  # Dark slate blue
        text_color = RGBColor(255, 255, 255)
    elif gender == "female":
        bg_color = RGBColor(192, 57, 43)  # Professional burgundy
        text_color = RGBColor(255, 255, 255)
    else:
        bg_color = RGBColor(41, 128, 185)  # Professional blue
        text_color = RGBColor(255, 255, 255)
    
    avatar.fill.solid()
    avatar.fill.fore_color.rgb = bg_color
    avatar.line.fill.background()
    
    # Add professional initials
    safe_name = safe_str(name)
    initials = ''.join([word[0].upper() for word in safe_name.split()[:2] if word])
    if not initials:
        initials = "P"
    
    text_frame = avatar.text_frame
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.text = initials
    text_frame.paragraphs[0].font.size = Pt(size * 20)
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].font.color.rgb = text_color
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return avatar

def add_clean_textbox(slide, left, top, width, height, text, font_size=11, 
                      bold=False, color=None, align=PP_ALIGN.LEFT):
    """Add a clean, professional textbox with None handling"""
    
    if color is None:
        color = RGBColor(51, 51, 51)
    
    # Convert None to empty string
    safe_text = remove_unicode(text)
    
    textbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.paragraphs[0].text = safe_text
    text_frame.paragraphs[0].font.size = Pt(font_size)
    text_frame.paragraphs[0].font.bold = bold
    text_frame.paragraphs[0].font.color.rgb = color
    text_frame.paragraphs[0].alignment = align
    
    return textbox

def add_section_header(slide, left, top, title):
    """Add a professional section header with underline"""
    
    safe_title = safe_str(title)
    
    # Title text
    title_box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(3), Inches(0.4)
    )
    title_frame = title_box.text_frame
    title_frame.text = safe_title
    title_frame.paragraphs[0].font.size = Pt(14)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 58, 147)
    
    # Underline
    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top + 0.35),
        Inches(0.8),
        Inches(0.03)
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = RGBColor(198, 155, 53)
    underline.line.fill.background()

def generate_ppt(data, filename):
    """Generate a professional, clean corporate resume PPT"""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Use blank slide layout
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Professional color palette
    COLOR_NAVY = RGBColor(31, 58, 147)      # Dark blue
    COLOR_GOLD = RGBColor(198, 155, 53)     # Gold accent
    COLOR_GRAY = RGBColor(128, 128, 128)    # Medium gray
    COLOR_LIGHT_GRAY = RGBColor(245, 245, 245)  # Very light gray
    COLOR_DARK = RGBColor(51, 51, 51)       # Dark text
    COLOR_WHITE = RGBColor(255, 255, 255)   # White
    
    # Clean white background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE
    
    # ===== LEFT COLUMN (Sidebar) =====
    # Light gray sidebar background
    sidebar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(3),
        Inches(7.5)
    )
    sidebar.fill.solid()
    sidebar.fill.fore_color.rgb = COLOR_LIGHT_GRAY
    sidebar.line.fill.background()
    
    # Avatar - Professional circular photo/initial
    create_professional_avatar(slide, 0.75, 0.4, 1.5, data.get("name", "Candidate"), data.get("gender", "neutral"))
    
    # Contact Information on sidebar
    contact_y = 2.2
    
    # Section: Contact
    add_clean_textbox(slide, 0.4, contact_y, 2.2, 0.3, "CONTACT", 
                     font_size=11, bold=True, color=COLOR_NAVY)
    
    # Phone
    phone = data.get('phone', '')
    if phone and phone != 'Not specified' and phone != 'Not found':
        add_clean_textbox(slide, 0.4, contact_y + 0.35, 2.2, 0.25, 
                         f"{safe_str(phone)}", font_size=8, color=COLOR_GRAY)
    
    # Email
    email = data.get('email', '')
    if email and email != 'Not specified' and email != 'Not found':
        add_clean_textbox(slide, 0.4, contact_y + 0.65, 2.2, 0.25, 
                         f"{safe_str(email)}", font_size=8, color=COLOR_GRAY)
    
    # Location
    location = data.get('location', '')
    if location and location != 'Not specified':
        add_clean_textbox(slide, 0.4, contact_y + 0.95, 2.2, 0.25, 
                         f"{safe_str(location)}", font_size=8, color=COLOR_GRAY)
    
    # LinkedIn
    linkedin = data.get('linkedin', '')
    if linkedin and linkedin != 'Not specified' and linkedin != 'Not found':
        add_clean_textbox(slide, 0.4, contact_y + 1.25, 2.2, 0.25, 
                         f"{safe_str(linkedin)[:20]}", font_size=8, color=COLOR_GRAY)
    
    # Skills Section
    skills_y = contact_y + 1.8
    add_clean_textbox(slide, 0.4, skills_y, 2.2, 0.3, "SKILLS", 
                     font_size=11, bold=True, color=COLOR_NAVY)
    
    skills = data.get('skills', [])
    skill_y = skills_y + 0.4
    for skill in skills[:8]:
        if skill:
            add_clean_textbox(slide, 0.5, skill_y, 2.1, 0.25, 
                             f"• {safe_str(skill)}", font_size=8, color=COLOR_DARK)
            skill_y += 0.3
        if skill_y > 7.0:
            break
    
    # Certifications Section
    certifications = data.get('certifications', [])
    if certifications and skill_y < 6.5:
        cert_y = skill_y + 0.2
        add_clean_textbox(slide, 0.4, cert_y, 2.2, 0.3, "CERTIFICATIONS", 
                         font_size=11, bold=True, color=COLOR_NAVY)
        
        cert_y += 0.4
        for cert in certifications[:3]:
            if cert:
                add_clean_textbox(slide, 0.5, cert_y, 2.1, 0.25, 
                                 f"✓ {safe_str(cert)}", font_size=7, color=COLOR_GRAY)
                cert_y += 0.25
    
    # ===== RIGHT COLUMN (Main Content) =====
    right_x = 3.3
    
    # Name
    name = safe_str(data.get('name', 'Professional Candidate'))
    name_box = slide.shapes.add_textbox(
        Inches(right_x), Inches(0.5), Inches(6.2), Inches(0.7)
    )
    name_frame = name_box.text_frame
    name_frame.text = name.upper()
    name_frame.paragraphs[0].font.size = Pt(24)
    name_frame.paragraphs[0].font.bold = True
    name_frame.paragraphs[0].font.color.rgb = COLOR_NAVY
    
    # Current Role
    role = safe_str(data.get('current_role', 'Professional'))
    role_box = slide.shapes.add_textbox(
        Inches(right_x), Inches(1.1), Inches(6.2), Inches(0.4)
    )
    role_frame = role_box.text_frame
    role_frame.text = role
    role_frame.paragraphs[0].font.size = Pt(14)
    role_frame.paragraphs[0].font.color.rgb = COLOR_GOLD
    
    # Divider line
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(right_x),
        Inches(1.55),
        Inches(6.2),
        Inches(0.02)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = COLOR_GOLD
    divider.line.fill.background()
    
    # Professional Summary Section
    summary_y = 1.8
    add_section_header(slide, right_x, summary_y, "PROFESSIONAL SUMMARY")
    
    summary = safe_str(data.get('professional_summary', ''))
    if len(summary) > 400:
        summary = summary[:400] + "..."
    
    add_clean_textbox(slide, right_x, summary_y + 0.45, 6.2, 0.9, 
                     summary, font_size=9, color=COLOR_GRAY)
    
    # Work Experience Section
    exp_y = 2.8
    add_section_header(slide, right_x, exp_y, "WORK EXPERIENCE")
    
    experiences = data.get('latest_3_experiences', [])
    exp_item_y = exp_y + 0.45
    
    for idx, exp in enumerate(experiences[:3]):
        if exp_item_y > 6.2:
            break
            
        # Company and Role
        company = safe_str(exp.get('company', 'Company Name'))
        role = safe_str(exp.get('role', 'Position'))
        
        company_role_text = f"{role}  |  {company}"
        add_clean_textbox(slide, right_x, exp_item_y, 5, 0.3, 
                         company_role_text, font_size=11, bold=True, color=COLOR_NAVY)
        
        # Duration
        duration = safe_str(exp.get('duration', 'Date Range'))
        add_clean_textbox(slide, right_x + 4.5, exp_item_y, 1.7, 0.3, 
                         duration, font_size=8, color=COLOR_GRAY, align=PP_ALIGN.RIGHT)
        
        exp_item_y += 0.35
        
        # Responsibilities
        responsibilities = exp.get('responsibilities', [])
        for resp in responsibilities[:3]:
            if resp:
                resp_text = safe_str(resp)
                if len(resp_text) > 70:
                    resp_text = resp_text[:67] + "..."
                add_clean_textbox(slide, right_x + 0.15, exp_item_y, 6, 0.25, 
                                 f"• {resp_text}", font_size=8, color=COLOR_DARK)
                exp_item_y += 0.25
        
        exp_item_y += 0.2
    
    # Education Section
    if exp_item_y < 6.0:
        edu_y = exp_item_y + 0.1
        add_section_header(slide, right_x, edu_y, "EDUCATION")
        
        education = data.get('education', {})
        if education:
            degree = safe_str(education.get('degree', 'Degree'))
            institution = safe_str(education.get('institution', 'Institution'))
            year = safe_str(education.get('year', 'Year'))
            
            edu_text = f"{degree}  |  {institution}"
            add_clean_textbox(slide, right_x, edu_y + 0.45, 5, 0.3, 
                             edu_text, font_size=10, bold=True, color=COLOR_NAVY)
            
            add_clean_textbox(slide, right_x + 4.5, edu_y + 0.45, 1.7, 0.3, 
                             year, font_size=8, color=COLOR_GRAY, align=PP_ALIGN.RIGHT)
    
    # Fit Score Badge (subtle, top right)
    fit_score = data.get('fit_score', 85)
    if fit_score is None:
        fit_score = 85
    
    # Score background (subtle circle)
    score_bg = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(8.5),
        Inches(0.4),
        Inches(1.2),
        Inches(1.2)
    )
    score_bg.fill.solid()
    
    if fit_score >= 80:
        score_bg.fill.fore_color.rgb = RGBColor(46, 204, 113)
    elif fit_score >= 60:
        score_bg.fill.fore_color.rgb = RGBColor(241, 196, 15)
    else:
        score_bg.fill.fore_color.rgb = RGBColor(231, 76, 60)
    
    score_bg.line.fill.background()
    
    # Score text
    score_text = score_bg.text_frame
    score_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    score_text.text = str(fit_score)
    score_text.paragraphs[0].font.size = Pt(22)
    score_text.paragraphs[0].font.bold = True
    score_text.paragraphs[0].font.color.rgb = COLOR_WHITE
    score_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Score label
    score_label = slide.shapes.add_textbox(
        Inches(8.5),
        Inches(1.65),
        Inches(1.2),
        Inches(0.25)
    )
    score_label.text_frame.text = "FIT SCORE"
    score_label.text_frame.paragraphs[0].font.size = Pt(7)
    score_label.text_frame.paragraphs[0].font.bold = True
    score_label.text_frame.paragraphs[0].font.color.rgb = COLOR_GRAY
    score_label.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Footer
    footer = slide.shapes.add_textbox(
        Inches(right_x),
        Inches(7.1),
        Inches(6.2),
        Inches(0.2)
    )
    footer.text_frame.text = f"AI Resume Intelligence Report • Generated {datetime.now().strftime('%B %d, %Y')}"
    footer.text_frame.paragraphs[0].font.size = Pt(6)
    footer.text_frame.paragraphs[0].font.color.rgb = COLOR_GRAY
    footer.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Save the presentation
    output_folder = "generated_ppts"
    os.makedirs(output_folder, exist_ok=True)
    output_path = os.path.join(output_folder, filename)
    prs.save(output_path)
    
    return output_path